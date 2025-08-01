import io
import os
from typing import Optional

from docx import Document
import openpyxl
from pptx import Presentation
import pandas as pd
from PIL import Image
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.utils import ImageReader

from fastapi import HTTPException
from loguru import logger

from ocr import process_vlm_file_with_retry, MemoryDataWriter

# Supported file extensions (migrated from ocr.py)
pdf_extensions = [".pdf"]
office_extensions = [".ppt", ".pptx", ".doc", ".docx", ".xls", ".xlsx"]
image_extensions = [".png", ".jpg", ".jpeg"]
text_extensions = [".txt", ".md", ".markdown", ".text", ".rst", ".log"]
all_supported_extensions = pdf_extensions + office_extensions + image_extensions + text_extensions


def detect_file_type(file_extension: str) -> str:
    """Detect file type (migrated from ocr.py)"""
    if file_extension in pdf_extensions:
        return "PDF"
    elif file_extension in office_extensions:
        return "Office"
    elif file_extension in image_extensions:
        return "Image"
    elif file_extension in text_extensions:
        return "Text"
    else:
        return "Unknown"


def _detect_office_file_type(file_bytes: bytes) -> str:
    """Detect specific Office file type by examining ZIP contents"""
    try:
        import zipfile
        import io
        
        # Create a ZIP file object from the bytes
        with zipfile.ZipFile(io.BytesIO(file_bytes), 'r') as zip_file:
            file_list = zip_file.namelist()
            
            # Check for Excel-specific files
            if any('xl/' in f for f in file_list) or any('worksheets/' in f for f in file_list):
                logger.info("Detected as Excel file (.xlsx)")
                return '.xlsx'
            
            # Check for PowerPoint-specific files
            elif any('ppt/' in f for f in file_list) or any('slides/' in f for f in file_list):
                logger.info("Detected as PowerPoint file (.pptx)")
                return '.pptx'
            
            # Check for Word-specific files
            elif any('word/' in f for f in file_list) or 'document.xml' in file_list:
                logger.info("Detected as Word file (.docx)")
                return '.docx'
            
            # Fallback: check content types
            elif '[Content_Types].xml' in file_list:
                try:
                    content_types = zip_file.read('[Content_Types].xml').decode('utf-8')
                    if 'spreadsheetml' in content_types:
                        logger.info("Detected as Excel file (.xlsx) via content types")
                        return '.xlsx'
                    elif 'presentationml' in content_types:
                        logger.info("Detected as PowerPoint file (.pptx) via content types")
                        return '.pptx'
                    elif 'wordprocessingml' in content_types:
                        logger.info("Detected as Word file (.docx) via content types")
                        return '.docx'
                except Exception as e:
                    logger.warning(f"Could not read content types: {e}")
            
            # Default fallback to docx if we can't determine
            logger.warning("Could not determine specific Office file type, defaulting to .docx")
            return '.docx'
            
    except Exception as e:
        logger.warning(f"Error detecting Office file type: {e}, defaulting to .docx")
        return '.docx'


def detect_file_type_from_content(file_bytes: bytes) -> str:
    """Detect file type from content (file signature) (migrated from ocr.py)"""
    if not file_bytes:
        logger.warning("Empty file_bytes provided to detect_file_type_from_content")
        return ""
    
    try:
        logger.info(f"Detecting file type from content. First 20 bytes: {file_bytes[:20]}")
        
        # PDF files
        if file_bytes.startswith(b'%PDF-'):
            logger.info("Detected as PDF file")
            return '.pdf'
        
        # Image files
        elif file_bytes.startswith(b'\x89PNG\r\n\x1a\n'):
            logger.info("Detected as PNG file")
            return '.png'
        elif file_bytes.startswith(b'\xff\xd8\xff'):
            logger.info("Detected as JPG file")
            return '.jpg'
        
        # Office files (ZIP-based)
        elif file_bytes.startswith(b'PK'):
            logger.info("Detected as Office file (ZIP-based), checking specific type...")
            return _detect_office_file_type(file_bytes)
        
        # Try to decode as text
        else:
            try:
                decoded = file_bytes[:1024].decode('utf-8')
                logger.info(f"Successfully decoded as UTF-8 text. Sample: {decoded[:50]}")
                return '.txt'
            except UnicodeDecodeError as e:
                logger.info(f"UTF-8 decode failed: {e}, trying GBK")
                try:
                    decoded = file_bytes[:1024].decode('gbk')
                    logger.info(f"Successfully decoded as GBK text. Sample: {decoded[:50]}")
                    return '.txt'
                except UnicodeDecodeError as e2:
                    logger.warning(f"Both UTF-8 and GBK decode failed: {e2}")
                    return ""
    except Exception as e:
        logger.exception(f"Error in detect_file_type_from_content: {e}")
        return ""


def check_document_complexity(file_bytes: bytes, file_extension: str) -> dict:
    """
    Check document complexity and recommend timeout based on file size (migrated from ocr.py)
    
    This function analyzes document characteristics to:
    - Estimate processing difficulty based on file size (most efficient method)
    - Recommend appropriate timeout values for VLM processing
    - Help optimize resource allocation for different document types
    """
    file_size = len(file_bytes)
    file_size_mb = file_size / (1024 * 1024)
    
    try:
        if file_extension in pdf_extensions:
            # Estimate complexity based on file size (faster than parsing PDF)
            if file_size_mb <= 2:
                recommended_timeout = 60
                complexity = "low"
            elif file_size_mb <= 10:
                recommended_timeout = 120
                complexity = "medium"
            elif file_size_mb <= 50:
                recommended_timeout = 300
                complexity = "high"
            else:
                recommended_timeout = 600
                complexity = "very_high"
            
            return {
                "file_size_mb": round(file_size_mb, 2),
                "complexity": complexity,
                "recommended_timeout": recommended_timeout
            }
        
        elif file_extension in image_extensions:
            if file_size_mb > 5:
                return {"file_size_mb": round(file_size_mb, 2), "complexity": "medium", "recommended_timeout": 120}
            else:
                return {"file_size_mb": round(file_size_mb, 2), "complexity": "low", "recommended_timeout": 60}
        
        elif file_extension in office_extensions:
            if file_size_mb > 10:
                return {"file_size_mb": round(file_size_mb, 2), "complexity": "high", "recommended_timeout": 180}
            else:
                return {"file_size_mb": round(file_size_mb, 2), "complexity": "medium", "recommended_timeout": 120}
        
        else:
            return {"file_size_mb": round(file_size_mb, 2), "complexity": "low", "recommended_timeout": 60}
            
    except Exception as e:
        logger.warning(f"Complexity check failed: {e}")
        return {"file_size_mb": round(file_size_mb, 2), "complexity": "medium", "recommended_timeout": 120}


async def process_text_and_office_documents(
    file_bytes: bytes, 
    file_extension: str, 
    timeout: int = None
) -> str:
    """处理文本文件和Office文档"""
    
    if file_extension in text_extensions:
        # 直接处理文本文件
        logger.info(f"Processing text file directly: {file_extension}")
        try:
            # 尝试UTF-8解码
            return file_bytes.decode('utf-8')
        except UnicodeDecodeError:
            try:
                # 回退到GBK编码
                return file_bytes.decode('gbk')
            except UnicodeDecodeError:
                try:
                    # 回退到Latin1编码
                    return file_bytes.decode('latin1')
                except UnicodeDecodeError:
                    logger.error("Failed to decode text file with any encoding")
                    raise HTTPException(
                        status_code=400,
                        detail="Cannot decode text file. File may be corrupted or in an unsupported encoding."
                    )
    
    elif file_extension in office_extensions:
        # 使用第三方库处理Office文档
        logger.info(f"Processing Office document with third-party library: {file_extension}")
        
        try:
            if file_extension in ['.docx', '.doc']:
                # 处理Word文档
                text_content = _extract_word_content(file_bytes)
                
            elif file_extension in ['.xlsx', '.xls']:
                # 处理Excel文档
                text_content = _extract_excel_content(file_bytes)
                
            elif file_extension in ['.pptx', '.ppt']:
                # 处理PowerPoint文档
                text_content = _extract_powerpoint_content(file_bytes)
                
            else:
                raise HTTPException(
                    status_code=400,
                    detail=f"Unsupported office document type: {file_extension}"
                )
            
            return text_content
            
        except Exception as e:
            logger.error(f"Failed to process office document {file_extension}: {str(e)}")
            raise HTTPException(
                status_code=500,
                detail=f"Failed to process office document: {str(e)}"
            )
    
    else:
        raise HTTPException(
            status_code=400,
            detail=f"Unsupported file type for text/office processing: {file_extension}"
        )


async def process_image_files(
    file_bytes: bytes, 
    file_extension: str, 
    timeout: int = None
) -> str:
    """处理图像文件，先转换为PDF再使用MinerU的VLM管线"""
    
    if file_extension not in image_extensions:
        raise HTTPException(
            status_code=400,
            detail=f"Unsupported image file type: {file_extension}"
        )
    
    logger.info(f"Processing image file by converting to PDF: {file_extension}")
    
    try:
        # 将图像转换为PDF
        pdf_bytes = _convert_image_to_pdf(file_bytes)
        
        # 计算复杂度信息
        complexity_info = check_document_complexity(pdf_bytes, '.pdf')
        
        # 使用VLM处理转换后的PDF
        image_writer = MemoryDataWriter()
        
        _, _, md_content = await process_vlm_file_with_retry(
            file_bytes=pdf_bytes,
            file_extension='.pdf',  # 转换后的格式
            image_writer=image_writer,
            backend=None,
            timeout=timeout,
            max_retries=3,
            complexity_info=complexity_info
        )
        
        return md_content if md_content else ""
        
    except Exception as e:
        logger.error(f"Failed to process image file {file_extension}: {str(e)}")
        raise HTTPException(
            status_code=500,
            detail=f"Failed to process image file: {str(e)}"
        )


async def process_pdf_files(
    file_bytes: bytes, 
    file_extension: str, 
    timeout: int = None
) -> str:
    """处理PDF文件，直接使用MinerU的VLM管线"""
    
    if file_extension not in pdf_extensions:
        raise HTTPException(
            status_code=400,
            detail=f"Unsupported PDF file type: {file_extension}"
        )
    
    logger.info(f"Processing PDF file with VLM: {file_extension}")
    
    try:
        # 计算复杂度信息
        complexity_info = check_document_complexity(file_bytes, file_extension)
        
        image_writer = MemoryDataWriter()
        
        _, _, md_content = await process_vlm_file_with_retry(
            file_bytes=file_bytes,
            file_extension=file_extension,
            image_writer=image_writer,
            backend=None,
            timeout=timeout,
            max_retries=3,
            complexity_info=complexity_info
        )
        
        return md_content if md_content else ""
        
    except Exception as e:
        logger.error(f"Failed to process PDF file {file_extension}: {str(e)}")
        raise HTTPException(
            status_code=500,
            detail=f"Failed to process PDF file: {str(e)}"
        )


def _extract_word_content(file_bytes: bytes) -> str:
    """从Word文档中提取文本内容"""
    try:
        doc = Document(io.BytesIO(file_bytes))
        full_text = []
        
        # 提取段落文本
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                full_text.append(paragraph.text)
        
        # 提取表格文本
        for table in doc.tables:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text.strip())
                if row_text:
                    full_text.append(' | '.join(row_text))
        
        return '\n'.join(full_text)
        
    except Exception as e:
        logger.error(f"Error extracting Word content: {str(e)}")
        raise


def _extract_excel_content(file_bytes: bytes) -> str:
    """从Excel文档中提取文本内容"""
    try:
        # 使用pandas读取Excel文件
        df_dict = pd.read_excel(io.BytesIO(file_bytes), sheet_name=None)
        
        full_text = []
        
        for sheet_name, df in df_dict.items():
            full_text.append(f"Sheet: {sheet_name}")
            
            # 转换DataFrame为文本
            if not df.empty:
                # 包含列名
                full_text.append(df.to_string(index=False))
            else:
                full_text.append("(Empty sheet)")
            
            full_text.append("")  # 空行分隔不同的sheet
        
        return '\n'.join(full_text)
        
    except Exception as e:
        logger.error(f"Error extracting Excel content: {str(e)}")
        raise


def _extract_powerpoint_content(file_bytes: bytes) -> str:
    """从PowerPoint文档中提取文本内容"""
    try:
        prs = Presentation(io.BytesIO(file_bytes))
        full_text = []
        
        for i, slide in enumerate(prs.slides, 1):
            full_text.append(f"Slide {i}:")
            
            # 提取幻灯片中的文本
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            
            if slide_text:
                full_text.extend(slide_text)
            else:
                full_text.append("(No text content)")
            
            full_text.append("")  # 空行分隔不同的幻灯片
        
        return '\n'.join(full_text)
        
    except Exception as e:
        logger.error(f"Error extracting PowerPoint content: {str(e)}")
        raise


def _convert_image_to_pdf(image_bytes: bytes) -> bytes:
    """将图像转换为PDF格式"""
    try:
        # 使用PIL打开图像
        image = Image.open(io.BytesIO(image_bytes))
        
        # 转换为RGB模式（如果需要）
        if image.mode != 'RGB':
            image = image.convert('RGB')
        
        # 创建PDF
        pdf_buffer = io.BytesIO()
        
        # 获取图像尺寸
        img_width, img_height = image.size
        
        # 计算PDF页面大小（保持图像比例）
        # 使用A4页面大小作为基准
        a4_width, a4_height = A4
        
        # 计算缩放比例
        scale_x = a4_width / img_width
        scale_y = a4_height / img_height
        scale = min(scale_x, scale_y)
        
        new_width = img_width * scale
        new_height = img_height * scale
        
        # 创建PDF文档
        c = canvas.Canvas(pdf_buffer, pagesize=(new_width, new_height))
        
        # 将图像添加到PDF
        img_buffer = io.BytesIO()
        image.save(img_buffer, format='PNG')
        img_buffer.seek(0)
        
        img_reader = ImageReader(img_buffer)
        c.drawImage(img_reader, 0, 0, width=new_width, height=new_height)
        c.save()
        
        pdf_buffer.seek(0)
        return pdf_buffer.getvalue()
        
    except Exception as e:
        logger.error(f"Error converting image to PDF: {str(e)}")
        raise HTTPException(
            status_code=500,
            detail=f"Failed to convert image to PDF: {str(e)}"
        )
